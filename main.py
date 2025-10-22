# from fastapi import FastAPI, Query
# from fastapi.responses import FileResponse, JSONResponse
# from pptx import Presentation
# from pptx.util import Inches, Pt
# import os

# app = FastAPI(title="QuickBase to PPT API")

# # Root endpoint
# @app.get("/")
# def root():
#     return {"message": "QuickBase to PPT API is running. Use /health to check status."}

# # Health check endpoint
# @app.get("/health")
# def health_check():
#     return {"status": "OK"}

# # Generate PPT endpoint
# @app.get("/generateppt")
# def generate_ppt(
#     templateId: str = Query(..., description="Quickbase Template ID"),
#     recordId: str = Query(..., description="Quickbase Record ID"),
#     filename: str = Query("output", description="Output PPT filename")
# ):
#     """
#     Generate a simple PPT for demonstration purposes.
#     """
#     try:
#         # Create a new presentation
#         prs = Presentation()

#         # Add a title slide
#         slide_layout = prs.slide_layouts[0]  # 0 = Title slide
#         slide = prs.slides.add_slide(slide_layout)
#         slide.shapes.title.text = f"QuickBase Record {recordId}"
#         slide.placeholders[1].text = f"Template ID: {templateId}"

#         # Add a content slide
#         slide_layout = prs.slide_layouts[1]  # 1 = Title + Content
#         slide = prs.slides.add_slide(slide_layout)
#         slide.shapes.title.text = "Record Details"
#         content = slide.placeholders[1]
#         content.text = f"This PPT was generated for Record ID {recordId} using Template ID {templateId}."

#         # Save the PPT
#         output_file = f"{filename}.pptx"
#         prs.save(output_file)

#         # Return the PPT as a downloadable file
#         return FileResponse(path=output_file, filename=output_file, media_type='application/vnd.openxmlformats-officedocument.presentationml.presentation')

#     except Exception as e:
#         return JSONResponse(content={"error": str(e)})


from fastapi import FastAPI, Query
from fastapi.responses import FileResponse, JSONResponse
from pptx import Presentation
from pptx.util import Inches, Pt
import json
import os

app = FastAPI(title="QuickBase to PPT API")

CONFIG_FILE = "config.json"

# Root endpoint
@app.get("/")
def root():
    return {"message": "QuickBase to PPT API is running. Use /health to check status."}

# Health check endpoint
@app.get("/health")
def health_check():
    return {"status": "OK"}

def add_table(slide, placeholder_name, table_data):
    """
    Add a table to the slide.
    For simplicity, we create a single-column table with each item as a row.
    """
    rows = len(table_data)
    cols = 1
    left = Inches(1)
    top = Inches(2)
    width = Inches(8)
    height = Inches(0.8 * rows)

    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    for i, row_text in enumerate(table_data):
        table.cell(i, 0).text = row_text
        for paragraph in table.cell(i, 0).text_frame.paragraphs:
            paragraph.font.size = Pt(12)

def add_bullets(slide, placeholder_name, bullets):
    """
    Add bullet points to the slide's content placeholder.
    """
    if len(slide.placeholders) > 1:
        content = slide.placeholders[1]
        content.text = ""
        for bullet in bullets:
            p = content.text_frame.add_paragraph()
            p.text = bullet
            p.level = 0
            p.font.size = Pt(12)

@app.get("/generateppt")
def generate_ppt(filename: str = Query("output", description="Output PPT filename")):
    try:
        # Load config.json
        if not os.path.exists(CONFIG_FILE):
            return JSONResponse(content={"error": f"{CONFIG_FILE} not found."}, status_code=400)
        
        with open(CONFIG_FILE, "r") as f:
            config = json.load(f)

        prs = Presentation()

        # Loop through slides
        for slide_config in config.get("slides", []):
            slide_number = slide_config.get("slide_number", 1)
            
            # Choose layout: 0 = Title, 1 = Title + Content
            slide_layout = prs.slide_layouts[0] if slide_number == 1 else prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)

            # Add placeholders
            placeholders = slide_config.get("placeholders", {})
            if placeholders:
                for key, value in placeholders.items():
                    # Replace text in title if placeholder is in title
                    if key.lower().find("title") != -1 or slide_number == 1:
                        slide.shapes.title.text = value
                    elif len(slide.placeholders) > 1:
                        slide.placeholders[1].text = value

            # Add tables
            tables = slide_config.get("tables", {})
            for key, table_data in tables.items():
                add_table(slide, key, table_data)

            # Add bullets
            bullets = slide_config.get("bullets", {})
            for key, bullet_list in bullets.items():
                add_bullets(slide, key, bullet_list)

        # Save PPT
        output_file = f"{filename}.pptx"
        prs.save(output_file)

        return FileResponse(
            path=output_file,
            filename=output_file,
            media_type='application/vnd.openxmlformats-officedocument.presentationml.presentation'
        )

    except Exception as e:
        return JSONResponse(content={"error": str(e)}, status_code=500)
