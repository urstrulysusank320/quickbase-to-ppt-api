# convert.py
import argparse
from pptx import Presentation
from pptx.util import Inches, Pt
from bs4 import BeautifulSoup
import os

def create_slide(prs, title, content_lines):
    """Add a slide with a title and bullet points/content."""
    slide_layout = prs.slide_layouts[1]  # Title + Content layout
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title

    body = slide.placeholders[1]  # Content placeholder
    for line in content_lines:
        p = body.text_frame.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(12)

def parse_html_to_slides(html_content):
    """
    Parses HTML content and returns a list of slides:
    Each slide is a tuple: (title, [content_lines])
    """
    soup = BeautifulSoup(html_content, 'lxml')
    slides = []

    current_title = "Title Slide"
    current_content = []

    # Iterate through all tags in body
    body = soup.body if soup.body else soup
    for tag in body.descendants:
        if tag.name in ['h1', 'h2', 'h3']:
            # Save previous slide if it has content
            if current_content:
                slides.append((current_title, current_content))
                current_content = []
            current_title = tag.get_text(strip=True)
        elif tag.name in ['p', 'li']:
            text = tag.get_text(strip=True)
            if text:
                current_content.append(text)
        elif tag.name == 'table':
            # Convert table to bullet lines
            for row in tag.find_all('tr'):
                row_text = " | ".join([cell.get_text(strip=True) for cell in row.find_all(['td', 'th'])])
                if row_text:
                    current_content.append(row_text)

    # Add last slide
    if current_content:
        slides.append((current_title, current_content))

    return slides

def generate_ppt(html_file, output_file):
    # Read HTML content
    if not os.path.exists(html_file):
        raise FileNotFoundError(f"{html_file} does not exist")
    
    with open(html_file, 'r', encoding='utf-8') as f:
        html_content = f.read()

    prs = Presentation()
    slides = parse_html_to_slides(html_content)

    for title, content in slides:
        create_slide(prs, title, content)

    prs.save(output_file)
    print(f"PPT generated: {output_file}")

if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Convert HTML Doc Template to PPT")
    parser.add_argument('--html_file', required=True, help='Input HTML file path')
    parser.add_argument('--output_name', required=True, help='Output PPT file name (without extension)')
    args = parser.parse_args()

    output_file = f"{args.output_name}.pptx"
    generate_ppt(args.html_file, output_file)
